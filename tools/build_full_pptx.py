"""
Build the full 22-slide DataLens AI end-semester PPTX deck.

Run from repo root:
    python tools/build_full_pptx.py
Output:
    reports/datalens_full_deck.pptx

This is a self-contained generator. It mirrors the editorial-dark design
language used by tools/build_pptx.py but follows the 22-slide outline the
team requested (Title → Contents → Intro → Problem → Motivation →
Literature → Gaps → Objectives → Architecture → … → Thank You).

It uses:
  • python-pptx for layout + native tables
  • One light hero band on the cover and section-divider slides
  • Editorial palette: cream + teal on near-black
  • JetBrains Mono for code/eyebrow strings, Inter elsewhere
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt


ROOT = Path(__file__).resolve().parent.parent

# ── Editorial colour palette ─────────────────────────────────────────────────

INK_1 = RGBColor(0x0A, 0x0A, 0x0A)
INK_2 = RGBColor(0x33, 0x33, 0x33)
INK_3 = RGBColor(0x55, 0x55, 0x55)
PAGE = RGBColor(0xFF, 0xFF, 0xFF)
BG_2 = RGBColor(0xF5, 0xF8, 0xFF)
BG_3 = RGBColor(0xF8, 0xFA, 0xFF)
BG_TABLE = RGBColor(0xFB, 0xFC, 0xFF)
ACCENT = RGBColor(0x0D, 0x8A, 0x72)   # editorial teal
ACCENT2 = RGBColor(0x25, 0x63, 0xEB)  # blue accent
WARN = RGBColor(0xC4, 0x7A, 0x09)
OK = RGBColor(0x2F, 0x7A, 0x4D)
BAD = RGBColor(0xB0, 0x37, 0x37)


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def add_textbox(slide, left, top, width, height, text, *,
                size=14, bold=False, color=INK_2, align=PP_ALIGN.LEFT,
                font="Inter", anchor=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *,
                size=14, color=INK_2, font="Inter", line_spacing=1.25,
                bullet_char="•"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"{bullet_char}  {item}" if bullet_char else item
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, color, *, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, left, top, width, height, fill, *, line=None, line_w=0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w if line_w else 0.75)
    shape.shadow.inherit = False
    return shape


def add_header(slide, prs, *, eyebrow="DataLens · AI", title="", subtitle=""):
    sw = prs.slide_width / 914400
    add_rect(slide, 0, 0, sw, 0.35, ACCENT)
    add_textbox(slide, 0.55, 0.45, 8, 0.3,
                eyebrow.upper(),
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    if title:
        add_textbox(slide, 0.55, 0.78, sw - 1.1, 0.85,
                    title, size=30, bold=True, color=INK_1)
    if subtitle:
        add_textbox(slide, 0.55, 1.55, sw - 1.1, 0.45,
                    subtitle, size=13, color=INK_3)
    add_rect(slide, 0.55, 2.05, 1.4, 0.04, ACCENT)


def add_footer(slide, prs, page_num: int, total: int):
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400
    add_rect(slide, 0.55, sh - 0.45, sw - 1.1, 0.01, INK_3)
    add_textbox(slide, 0.55, sh - 0.4, 5, 0.3,
                "DataLens AI  ·  End-Semester Review  ·  Group 13",
                size=9, color=INK_3, font="JetBrains Mono")
    add_textbox(slide, sw - 1.55, sh - 0.4, 1, 0.3,
                f"{page_num:02d} / {total:02d}",
                size=9, color=INK_2, bold=True, align=PP_ALIGN.RIGHT,
                font="JetBrains Mono")


def style_cell(cell, text, *, size=10, bold=False, color=INK_2,
                font="Inter", fill=None, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.vertical_anchor = anchor
    tf.word_wrap = True
    # Clear existing
    p = tf.paragraphs[0]
    p.alignment = align
    if p.runs:
        for r in list(p.runs):
            r.text = ""
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_table(slide, *, left, top, width, height, headers, rows,
              header_fill=ACCENT, header_color=PAGE,
              row_fill=PAGE, alt_fill=BG_TABLE,
              header_size=11, body_size=10, first_col_bold=False,
              col_widths=None):
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
    )
    tbl = table_shape.table

    if col_widths is not None:
        total_specified = sum(col_widths)
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(width * (w / total_specified))

    # Header row
    for c, label in enumerate(headers):
        style_cell(tbl.cell(0, c), label,
                   size=header_size, bold=True, color=header_color,
                   fill=header_fill, font="Inter")

    # Body rows
    for r, row in enumerate(rows, start=1):
        fill = row_fill if (r % 2) else alt_fill
        for c, value in enumerate(row):
            style_cell(
                tbl.cell(r, c), value,
                size=body_size,
                bold=(first_col_bold and c == 0),
                color=INK_1 if (first_col_bold and c == 0) else INK_2,
                fill=fill, font="Inter",
            )
    return tbl


# ── Slide builders ───────────────────────────────────────────────────────────


def slide_01_title(prs):
    s = add_blank_slide(prs)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400

    # Hero band
    add_rect(s, 0, 0, sw, sh * 0.55, BG_3)
    add_rect(s, 0, sh * 0.55 - 0.04, sw, 0.04, ACCENT)

    add_textbox(s, 0.6, 0.7, 8, 0.4,
                "DATALENS · AI   ·   END-SEMESTER REVIEW",
                size=12, bold=True, color=ACCENT, font="JetBrains Mono")

    add_textbox(s, 0.6, 1.3, sw - 1.2, 1.4,
                "DataLens AI", size=64, bold=True, color=INK_1)

    add_textbox(s, 0.6, 2.85, sw - 1.2, 0.7,
                "Read the signal in your data.",
                size=24, bold=True, color=ACCENT)

    add_textbox(s, 0.6, 3.55, sw - 1.2, 0.6,
                "Offline AI-Powered Dataset Analyzer",
                size=14, color=INK_3)

    # Bottom info card
    card_top = sh - 2.95
    add_rect(s, 0.6, card_top, sw - 1.2, 2.25, BG_2)

    add_textbox(s, 0.85, card_top + 0.2, 5, 0.3, "GUIDE",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    add_textbox(s, 0.85, card_top + 0.5, 5, 0.4,
                "Prof. Kishor Pathak", size=14, bold=True, color=INK_1)

    add_textbox(s, 0.85, card_top + 1.0, 5, 0.3,
                "TEAM  ·  GROUP 13",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    add_bullets(s, 0.85, card_top + 1.3, 6, 1.1, [
        "Parth Birari   ·  53 / IT D / 1251130415",
        "Parth Dongre   ·  54 / IT D / 1251130421",
        "Aryan Patil    ·  55 / IT D / 1251130411",
        "Atharv Patil   ·  56 / IT D / 1251130441",
    ], size=12, color=INK_2)

    add_textbox(s, sw - 5.6, card_top + 0.2, 5, 0.3,
                "INSTITUTE",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono",
                align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, card_top + 0.5, 5, 0.4,
                "Vishwakarma Institute of Technology, Pune",
                size=14, bold=True, color=INK_1, align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, card_top + 0.95, 5, 0.4,
                "Department of Information Technology",
                size=12, color=INK_2, align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, card_top + 1.55, 5, 0.4,
                "Date:  May 2026",
                size=11, color=INK_3, align=PP_ALIGN.RIGHT,
                font="JetBrains Mono")
    return s


def slide_02_contents(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Contents",
               subtitle="Twenty-two slides — from problem framing to thank-you.")
    sw = prs.slide_width / 914400

    items = [
        "1.  Title",                           "2.  Contents",
        "3.  Introduction",                    "4.  Problem Statement",
        "5.  Motivation",                      "6.  Literature Review",
        "7.  Research Gaps",                   "8.  Objectives",
        "9.  System Architecture",             "10. Data Processing & Quality",
        "11. Statistical Analysis",            "12. Anomaly Detection Engine",
        "13. Machine Learning Pipeline",       "14. Explainable AI (SHAP)",
        "15. Time-Series & Drift Detection",   "16. Reporting & Visualization",
        "17. Agentic AI Layer",                "18. Tool Comparison",
        "19. Limitations & Future Work",       "20. Conclusion",
        "21. References",                      "22. Thank You",
    ]
    col_w = (sw - 1.3) / 2
    left_col = items[0::2]
    right_col = items[1::2]

    add_bullets(s, 0.55, 2.5, col_w, 4.6, left_col,
                size=13, color=INK_2, bullet_char="")
    add_bullets(s, 0.55 + col_w + 0.2, 2.5, col_w, 4.6, right_col,
                size=13, color=INK_2, bullet_char="")
    return s


def slide_03_introduction(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Introduction",
               subtitle="What is DataLens AI?")
    sw = prs.slide_width / 914400

    # Definition card
    add_rect(s, 0.55, 2.4, sw - 1.1, 1.3, BG_2)
    add_rect(s, 0.55, 2.4, sw - 1.1, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.55, sw - 1.5, 0.4,
                "DataLens AI is an offline-first analytics platform that:",
                size=13, bold=True, color=INK_1)
    add_bullets(s, 0.85, 2.95, sw - 1.6, 0.8, [
        "analyzes datasets automatically  ·  trains ML leaderboards  ·  detects anomalies and drift",
        "generates explainable AI insights  ·  supports conversational dataset Q&A — all locally.",
    ], size=12, color=INK_2)

    # Three column cards
    col_top = 4.0
    col_h = 2.7
    col_w = (sw - 1.5) / 3
    cols = [
        ("SUPPORTED FILE TYPES", ACCENT, [
            "CSV  ·  TSV",
            "Excel (.xlsx / .xls)",
            "JSON  ·  JSONL",
            "Encoding waterfall handles BOM, mojibake, latin-1.",
        ]),
        ("CORE FEATURES", ACCENT2, [
            "Automated profiling",
            "ML leaderboard generation",
            "SHAP global + per-row explanations",
            "Time-series analytics + drift detection",
            "AI-powered dataset assistant",
        ]),
        ("KEY PROPERTY", WARN, [
            "Mathematically grounded",
            "Explainable & reproducible",
            "Privacy-friendly (offline-first)",
            "Every claim traces to a named formula.",
        ]),
    ]
    for i, (title, accent, bullets) in enumerate(cols):
        x = 0.55 + i * (col_w + 0.2)
        add_rect(s, x, col_top, col_w, col_h, BG_2)
        add_rect(s, x, col_top, col_w, 0.06, accent)
        add_textbox(s, x + 0.2, col_top + 0.18, col_w - 0.4, 0.3,
                    title, size=10, bold=True, color=accent,
                    font="JetBrains Mono")
        add_bullets(s, x + 0.2, col_top + 0.55, col_w - 0.4, col_h - 0.6,
                    bullets, size=12, color=INK_2)
    return s


def slide_04_problem(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Problem Statement",
               subtitle="Where existing tooling falls short.")
    sw = prs.slide_width / 914400
    col_w = (sw - 1.3) / 2

    # Left card — fragmented tooling + privacy
    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, BAD)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Problems with Existing Systems",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, col_w - 0.4, 1.9, [
        "Profiling tools lack ML capabilities",
        "AutoML platforms ignore data quality",
        "Drift systems require deployed models",
        "Conversational AI hallucinates without grounding",
    ], size=12)
    add_textbox(s, 0.75, 5.1, col_w - 0.4, 0.4,
                "Privacy Risks",
                size=13, bold=True, color=BAD)
    add_bullets(s, 0.75, 5.45, col_w - 0.4, 1.4, [
        "Modern AI analytics tools are cloud-dependent,",
        "closed-source, and unsafe for sensitive datasets.",
    ], size=12)

    # Right card — real-world dataset challenges + goal
    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Real-world Dataset Challenges",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.15, col_w - 0.4, 2.6, [
        "Missing values  ·  mixed dtypes",
        "Encoding mismatches (BOM, latin-1)",
        "High-cardinality / ID-like columns",
        "Hidden anomalies  ·  distribution drift",
        "Irregular timestamps  ·  unicode pitfalls",
    ], size=12)

    # Goal banner
    add_rect(s, rx + 0.2, 5.7, col_w - 0.4, 1.0, BG_3)
    add_textbox(s, rx + 0.4, 5.78, col_w - 0.8, 0.3,
                "GOAL",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    add_textbox(s, rx + 0.4, 6.05, col_w - 0.8, 0.6,
                "Build one unified offline AI-powered analytics platform.",
                size=13, bold=True, color=INK_1)
    return s


def slide_05_motivation(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Motivation",
               subtitle="Why DataLens AI?")
    sw = prs.slide_width / 914400

    # Two columns
    col_w = (sw - 1.3) / 2

    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, INK_3)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Today's analytics workflow",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.2, col_w - 0.4, 3.5, [
        "Multiple disconnected tools",
        "Manual preprocessing per dataset",
        "Repetitive notebook scripting",
        "Cloud lock-in for AI insights",
        "Demands technical expertise from every analyst",
    ], size=13)

    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Our vision — one unified pipeline",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.2, col_w - 0.4, 3.5, [
        "Profiles data automatically",
        "Performs ML analysis with a leaderboard",
        "Explains predictions with SHAP",
        "Detects anomalies & drift",
        "Tracks time-series structure",
        "Supports natural-language Q&A — locally.",
    ], size=13)
    return s


def slide_06_literature(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Literature Review",
               subtitle="Prior work across explainable AI · table understanding · drift · anomaly · AutoML.")

    headers = ["Research Area", "Paper / System", "Core Idea",
               "Key Contribution", "Gap Identified"]

    rows = [
        ("Explainable AI", "SHAP — Lundberg & Lee (2017)",
         "Explain predictions mathematically",
         "Standard for feature attribution",
         "High compute on large models"),
        ("Conversational Analytics", "PandasAI",
         "NL interaction with datasets",
         "LLM-generated pandas / SQL execution",
         "Cloud dependence + hallucination"),
        ("Table Understanding (LLM)", "TableGPT",
         "Unified table reasoning + commands",
         "Querying, viz, prediction in one model",
         "Cloud-scale compute required"),
        ("Tabular LLMs (Office)", "TableLLM",
         "Spreadsheet manipulation by LLMs",
         "Distant supervision + code generation",
         "Limited ML diagnostics"),
        ("Few-Shot Tabular ML", "TabLLM",
         "LLM-driven tabular classification",
         "Competitive with classical tabular ML",
         "Prediction-only; no profiling / anomalies"),
        ("Table-Tuned LLMs", "Table-GPT",
         "LLM fine-tuned on diverse table tasks",
         "Better generalization on table reasoning",
         "Still depends on huge cloud LLMs"),
        ("Drift Detection", "Evidently AI",
         "Detect dataset & model drift",
         "Popularized production ML monitoring",
         "Needs deployed pipelines"),
        ("Anomaly Detection", "Isolation Forest (Liu et al. 2008)",
         "Anomalies via random partitioning",
         "Fast unsupervised detector",
         "Single-model misses contextual outliers"),
        ("Time-Series", "STL — Cleveland et al. (1990)",
         "Seasonal-trend decomposition (Loess)",
         "Strong trend / seasonal analysis",
         "Needs forecasting complement"),
        ("AutoML", "H2O.ai",
         "Automate ML training & selection",
         "Simplified ML development",
         "Weak profiling + explainability"),
    ]

    add_table(
        s,
        left=0.4, top=2.35, width=12.5, height=4.65,
        headers=headers, rows=rows,
        header_size=10, body_size=8.5,
        col_widths=[2.0, 2.7, 2.7, 2.8, 2.6],
        first_col_bold=True,
    )
    return s


def slide_07_gaps(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Research Gaps Identified",
               subtitle="What no single tool delivers — and how DataLens AI fills the gap.")
    sw = prs.slide_width / 914400
    col_w = (sw - 1.3) / 2

    # Existing limitations
    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, BAD)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Existing limitations",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.2, col_w - 0.4, 3.5, [
        "No unified offline analytics ecosystem",
        "Limited explainability in AutoML platforms",
        "Hallucination risk in conversational AI",
        "Weak integration of ML + profiling + drift",
        "Poor privacy guarantees on cloud assistants",
        "Specialised tools — each solves one domain only",
    ], size=12)

    # DataLens fills this by combining
    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "DataLens AI combines",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.2, col_w - 0.4, 3.5, [
        "Offline analytics  ·  privacy by default",
        "Explainable ML pipelines (SHAP + permutation)",
        "Multi-detector anomaly ensemble",
        "PSI / KS / chi-square drift analysis",
        "Time-series analytics (ADF · KPSS · STL · HW)",
        "RAG-powered local AI assistant (Ollama)",
    ], size=12)
    return s


def slide_08_objectives(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Objectives",
               subtitle="Seven things the system must do, end-to-end.")
    sw = prs.slide_width / 914400

    # Number tiles
    tiles = [
        ("01", "Auto-detect column roles, dtypes, and dataset signals.", ACCENT),
        ("02", "Produce 0–100 Health and ML-Readiness scores.", ACCENT2),
        ("03", "Run a cross-validated 7-model leaderboard for any target.", WARN),
        ("04", "Generate global + per-row SHAP attributions for the winner.", ACCENT),
        ("05", "Detect time-series structure (ADF + KPSS + STL + Holt-Winters).", ACCENT2),
        ("06", "Compute PSI / KS / chi-square drift between halves of one upload.", WARN),
        ("07", "Expose chat-style Q&A grounded in the analysis result via RAG.", ACCENT),
    ]
    top = 2.45
    tile_h = 0.62
    tile_w = sw - 1.1
    for i, (num, text, accent) in enumerate(tiles):
        y = top + i * (tile_h + 0.13)
        add_rect(s, 0.55, y, tile_w, tile_h, BG_2)
        add_rect(s, 0.55, y, 0.08, tile_h, accent)
        add_textbox(s, 0.75, y + 0.07, 0.7, 0.5,
                    num, size=20, bold=True, color=accent,
                    font="JetBrains Mono")
        add_textbox(s, 1.5, y + 0.13, tile_w - 1.1, 0.5,
                    text, size=13, color=INK_1)
    return s


def slide_09_architecture(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="System Architecture",
               subtitle="Six-phase pipeline · 40+ modules · parallel where safe.")
    sw = prs.slide_width / 914400

    # Six phase blocks across the slide
    phases = [
        ("P1", "INGEST",
         "load → profile → roles → signals → selector",
         ACCENT),
        ("P2", "QUALITY",
         "health → ML-readiness → auto-cleaner",
         ACCENT2),
        ("P3", "PARALLEL",
         "deep stats · anomalies · time-series · text",
         WARN),
        ("P4", "ML CHAIN",
         "leaderboard → SHAP explainability",
         RGBColor(0xA7, 0x8B, 0xFA)),
        ("P5", "REPORTING",
         "chart planner → visualizer → PDF builder",
         RGBColor(0x7D, 0xD3, 0xFC)),
        ("P6", "AGENTIC AI",
         "brief → RAG index → local LLM → answer",
         BAD),
    ]
    box_w = (sw - 1.1 - 0.15 * 5) / 6
    box_h = 2.55
    top = 2.55
    for i, (tag, name, blurb, color) in enumerate(phases):
        x = 0.55 + i * (box_w + 0.15)
        add_rect(s, x, top, box_w, box_h, BG_2)
        add_rect(s, x, top, box_w, 0.06, color)
        # Tag chip
        add_rect(s, x + 0.1, top + 0.18, 0.55, 0.32, color)
        add_textbox(s, x + 0.1, top + 0.22, 0.55, 0.3,
                    tag, size=11, bold=True, color=PAGE,
                    align=PP_ALIGN.CENTER, font="JetBrains Mono")
        add_textbox(s, x + 0.15, top + 0.65, box_w - 0.3, 0.4,
                    name, size=12, bold=True, color=INK_1)
        add_textbox(s, x + 0.15, top + 1.05, box_w - 0.3, 1.4,
                    blurb, size=10, color=INK_3, font="JetBrains Mono")

    # Surface labels above and artefact labels below
    add_textbox(s, 0.55, 2.25, sw - 1.1, 0.25,
                "React Dashboard (17 tabs)  ·  Streamlit Console  ·  Flask API :5055",
                size=10, color=INK_3, font="JetBrains Mono",
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.55, top + box_h + 0.15, sw - 1.1, 0.25,
                "JSON contract  ·  20+ chart PNGs  ·  cleaned CSV  ·  multi-page PDF report",
                size=10, color=INK_3, font="JetBrains Mono",
                align=PP_ALIGN.CENTER)

    # Safety footer
    add_textbox(s, 0.55, top + box_h + 0.55, sw - 1.1, 0.4,
                "_safe_call() wraps every parallel branch — single failure ⇒ "
                "{ \"available\": false } instead of sinking the run.",
                size=10, color=INK_3, font="JetBrains Mono",
                align=PP_ALIGN.CENTER)
    return s


def slide_10_data_processing(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Data Processing & Quality",
               subtitle="Ingest mechanics · weighted health · auto-cleaner.")
    sw = prs.slide_width / 914400

    # Three columns
    col_w = (sw - 1.5) / 3
    top = 2.5
    col_h = 4.4

    cards = [
        ("DATA INGESTION", ACCENT, [
            "Multi-format support  (CSV · TSV · Excel · JSON)",
            "Encoding waterfall:  utf-8-sig → utf-8 → windows-1252 → latin-1",
            "Secure CSV handling — werkzeug.secure_filename",
            "CSV-injection guard prefixes  =  +  -  @  cells with  '",
        ]),
        ("HEALTH SCORE", ACCENT2, [
            "Components — completeness · consistency · uniqueness · outlier safety",
            "Health = 0.30·Completeness + 0.20·Consistency + 0.20·Uniqueness + 0.20·Safety + 10",
            "Labels — ≥90 Excellent · ≥75 Good · ≥60 Moderate · ≥40 Poor · else Critical",
            "Outlier detection via 1.5·IQR per numeric column",
        ]),
        ("AUTO-CLEANING", WARN, [
            "Numeric NaN → median impute",
            "Categorical NaN → mode impute",
            "drop_duplicates() on exact-match rows",
            "Cleaned CSV → cleaned/<id>_cleaned.csv",
            "Health re-computed before vs after",
        ]),
    ]
    for i, (title, accent, bullets) in enumerate(cards):
        x = 0.55 + i * (col_w + 0.25)
        add_rect(s, x, top, col_w, col_h, BG_2)
        add_rect(s, x, top, col_w, 0.06, accent)
        add_textbox(s, x + 0.2, top + 0.18, col_w - 0.4, 0.3,
                    title, size=10, bold=True, color=accent,
                    font="JetBrains Mono")
        add_bullets(s, x + 0.2, top + 0.55, col_w - 0.4, col_h - 0.7,
                    bullets, size=11, color=INK_2)
    return s


def slide_11_statistics(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Statistical Analysis",
               subtitle="Univariate · bivariate · best-fit distribution.")
    sw = prs.slide_width / 914400

    col_w = (sw - 1.3) / 2

    # Left card — modules
    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Statistical modules",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, col_w - 0.4, 3.5, [
        "Per-column descriptives (mean, σ, IQR)",
        "Skewness γ₁ + label · Kurtosis γ₂ + label",
        "Three-view outlier flags (IQR · z · MAD)",
        "Normality battery — Shapiro · D'Agostino · Anderson",
        "AIC distribution fit across 6 candidates",
        "Bootstrap 95 % CIs for mean and median (BCa)",
    ], size=12)

    # Right card — bivariate
    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT2)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Bivariate tests",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.15, col_w - 0.4, 3.5, [
        "Numeric × numeric — Pearson · Spearman · Kendall",
        "Categorical × categorical — chi-square + Cramér's V",
        "Binary × numeric — point-biserial correlation",
        "Group differences — Mann-Whitney U / Kruskal-Wallis H",
        "Multicollinearity — VIFⱼ = 1 / (1 − R²ⱼ)",
        "Target leakage — flag |Pearson r| ≥ 0.98",
    ], size=12)

    return s


def slide_12_anomaly(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Anomaly Detection Engine",
               subtitle="Five normalised detectors averaged into a [0,1] consensus.")
    sw = prs.slide_width / 914400

    # Top — algorithm tiles
    tiles = [
        ("Isolation Forest",
         "score(x) = 2^(−E[h(x)]/c(n))",
         "shorter average path ⇒ more anomalous", ACCENT),
        ("Local Outlier Factor",
         "LOF_k(x) = mean( lrd(y) / lrd(x) )",
         "ratio of local density to k neighbours", ACCENT2),
        ("Robust z (MAD)",
         "z* = (x − median) / (1.4826·MAD)",
         "flag |z*| > 3.5  ·  breakdown-resilient", WARN),
        ("Mahalanobis",
         "D²(x) = (x−μ)ᵀ Σ⁻¹ (x−μ)",
         "Σ via Minimum-Covariance-Determinant", BAD),
    ]
    tile_w = (sw - 1.1 - 0.2 * 3) / 4
    tile_h = 1.6
    top = 2.5
    for i, (name, formula, blurb, color) in enumerate(tiles):
        x = 0.55 + i * (tile_w + 0.2)
        add_rect(s, x, top, tile_w, tile_h, BG_2)
        add_rect(s, x, top, tile_w, 0.06, color)
        add_textbox(s, x + 0.18, top + 0.15, tile_w - 0.36, 0.35,
                    name, size=12, bold=True, color=INK_1)
        add_textbox(s, x + 0.18, top + 0.55, tile_w - 0.36, 0.4,
                    formula, size=10, color=color,
                    font="JetBrains Mono")
        add_textbox(s, x + 0.18, top + 1.0, tile_w - 0.36, 0.5,
                    blurb, size=10, color=INK_3)

    # Bottom — ensemble explainer
    add_rect(s, 0.55, 4.4, sw - 1.1, 2.4, BG_3)
    add_rect(s, 0.55, 4.4, sw - 1.1, 0.06, ACCENT)
    add_textbox(s, 0.75, 4.55, sw - 1.5, 0.4,
                "Why an ensemble?",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 5.0, sw - 1.5, 1.7, [
        "Each detector is min-max normalised to [0, 1].",
        "Ensemble s(x) = mean across available detectors  ·  flag if s(x) ≥ 0.6.",
        "Improves robustness, lowers false positives, captures different anomaly patterns.",
        "Output  →  per-row scores + top-25 flagged rows + per-detector breakdown.",
    ], size=12)
    return s


def slide_13_ml_pipeline(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Machine Learning Pipeline",
               subtitle="Auto-detect task · 7-model leaderboard · 5-fold CV.")
    sw = prs.slide_width / 914400

    col_w = (sw - 1.3) / 2

    # Left — workflow
    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "ML workflow",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, col_w - 0.4, 3.5, [
        "Auto task detection — classification when target has ≤ 10 unique values; else regression.",
        "Unified preprocessing — median impute · stringify · OneHotEncoder(handle_unknown='ignore').",
        "5-fold StratifiedKFold (clf) · KFold (reg) · cross_validate.",
        "Winner refit on 75 % · validated on 25 % holdout.",
        "Diagnostics — confusion matrix · residual summary · top-10 worst predictions.",
    ], size=12)

    # Right — models + metrics
    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT2)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Models on the leaderboard",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.15, col_w - 0.4, 1.6, [
        "Random Forest  ·  Gradient Boosting  ·  XGBoost  ·  LightGBM",
        "Logistic / Linear (Ridge / Lasso)  ·  K-Neighbors  ·  Dummy baseline",
    ], size=12)

    add_textbox(s, rx + 0.2, 4.85, col_w - 0.4, 0.4,
                "Metrics reported",
                size=13, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 5.2, col_w - 0.4, 1.7, [
        "Classification — Accuracy · Precision · Recall · F1 weighted · Brier",
        "Regression — R² · MAE · RMSE",
        "Each metric: per-fold mean ± std",
    ], size=12)
    return s


def slide_14_shap(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Explainable AI (SHAP)",
               subtitle="Trace every prediction back to its features.")
    sw = prs.slide_width / 914400

    # Formula card
    add_rect(s, 0.55, 2.4, sw - 1.1, 1.05, BG_3)
    add_rect(s, 0.55, 2.4, sw - 1.1, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.55, sw - 1.5, 0.35,
                "Shapley value (Lundberg & Lee 2017)",
                size=12, bold=True, color=ACCENT,
                font="JetBrains Mono")
    add_textbox(s, 0.75, 2.85, sw - 1.5, 0.55,
                "φᵢ(x)  =  Σ_{S⊆F\\{i}}  |S|! (|F|−|S|−1)! / |F|!  ·  [v(S∪{i}) − v(S)]",
                size=14, bold=True, color=INK_1, font="JetBrains Mono")

    # Three sub-cards
    col_w = (sw - 1.5) / 3
    top = 3.7
    col_h = 3.2
    cards = [
        ("WHY EXPLAINABILITY", ACCENT, [
            "Most AI systems behave like black boxes.",
            "DataLens AI traces every prediction.",
            "Local accuracy:  Σᵢ φᵢ(x) + E[ŷ] = ŷ(x).",
        ]),
        ("OUTPUTS", ACCENT2, [
            "Top-15 global features (one-hot dummies summed back to parents)",
            "Top-3 per-row stories (e.g. 'price +$45.2K due to Location')",
            "SHAP beeswarm summary PNG saved to static/charts/",
        ]),
        ("IMPLEMENTATIONS", WARN, [
            "shap.TreeExplainer  →  RF · GB · XGB · LGBM (poly-time exact)",
            "shap.LinearExplainer  →  LogReg · Ridge · Lasso (closed form)",
            "Permutation importance — always-available cross-check",
        ]),
    ]
    for i, (title, accent, bullets) in enumerate(cards):
        x = 0.55 + i * (col_w + 0.25)
        add_rect(s, x, top, col_w, col_h, BG_2)
        add_rect(s, x, top, col_w, 0.06, accent)
        add_textbox(s, x + 0.18, top + 0.18, col_w - 0.36, 0.3,
                    title, size=10, bold=True, color=accent,
                    font="JetBrains Mono")
        add_bullets(s, x + 0.18, top + 0.55, col_w - 0.36, col_h - 0.7,
                    bullets, size=11, color=INK_2)
    return s


def slide_15_timeseries_drift(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Time-Series & Drift Detection",
               subtitle="ADF · KPSS · STL · Holt-Winters · PSI · KS · chi-square.")
    sw = prs.slide_width / 914400

    col_w = (sw - 1.3) / 2

    # Left — time-series
    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Time-series features",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, col_w - 0.4, 3.5, [
        "Auto date detection  —  parse-rate × monotonicity × name keyword",
        "Frequency inference + FFT-based period detection",
        "Stationarity:  ADF (H₀ unit root) + KPSS (H₀ stationary) → 4-way verdict",
        "STL decomposition  ·  trend / seasonal strengths via Hyndman 2018",
        "Forecast: Holt-Winters additive vs naive last-value baseline (MAE / RMSE)",
    ], size=12)

    # Right — drift
    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, WARN)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Drift detection",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.15, col_w - 0.4, 1.7, [
        "PSI  =  Σ ( p_cur − p_ref ) · ln( p_cur / p_ref )",
        "KS two-sample  ·  flag if p < 0.01",
        "Categorical chi-square independence test",
    ], size=12)

    add_textbox(s, rx + 0.2, 5.0, col_w - 0.4, 0.4,
                "Severity buckets (PSI)",
                size=13, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 5.35, col_w - 0.4, 1.6, [
        "< 0.10  →  stable",
        "< 0.25  →  minor",
        "< 0.50  →  moderate",
        "≥ 0.50  →  severe",
    ], size=12)
    return s


def slide_16_reporting(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Reporting & Visualization",
               subtitle="Role-aware chart planner · 20+ chart types · multi-page PDF.")
    sw = prs.slide_width / 914400

    # Three-column cards
    col_w = (sw - 1.5) / 3
    top = 2.5
    col_h = 4.4
    cards = [
        ("AUTOMATED REPORTING", ACCENT, [
            "Interactive React dashboard — 17 tabs",
            "Statistical & ML charts auto-rendered",
            "Multi-page styled PDF report",
            "Cleaned CSV download",
        ]),
        ("SMART CHART PLANNER", ACCENT2, [
            "Selects chart specs from column roles + dataset signals",
            "Skips meaningless charts (no histogram of an ID column)",
            "Light-up triggers for time-series, leaderboard, SHAP",
        ]),
        ("RENDERING STACK", WARN, [
            "Matplotlib (Agg backend) · Seaborn",
            "Plotly  ·  Highcharts (in React)",
            "Editorial dark palette mirrors the dashboard tokens",
            "PNGs saved to static/charts/<dataset_id>_*.png",
        ]),
    ]
    for i, (title, accent, bullets) in enumerate(cards):
        x = 0.55 + i * (col_w + 0.25)
        add_rect(s, x, top, col_w, col_h, BG_2)
        add_rect(s, x, top, col_w, 0.06, accent)
        add_textbox(s, x + 0.2, top + 0.18, col_w - 0.4, 0.3,
                    title, size=10, bold=True, color=accent,
                    font="JetBrains Mono")
        add_bullets(s, x + 0.2, top + 0.55, col_w - 0.4, col_h - 0.7,
                    bullets, size=11, color=INK_2)
    return s


def slide_17_agentic(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Agentic AI Layer",
               subtitle="Local LLM grounded by RAG · sandboxed pandas execution.")
    sw = prs.slide_width / 914400

    # Workflow strip
    steps = [
        ("Dataset", ACCENT),
        ("Embeddings", ACCENT2),
        ("RAG Top-k", WARN),
        ("LLM Writer", RGBColor(0xA7, 0x8B, 0xFA)),
        ("Answer", ACCENT),
    ]
    box_w = (sw - 1.1 - 0.4 * 4) / 5
    box_h = 0.9
    y = 2.5
    for i, (name, color) in enumerate(steps):
        x = 0.55 + i * (box_w + 0.4)
        add_rect(s, x, y, box_w, box_h, BG_2)
        add_rect(s, x, y, box_w, 0.06, color)
        add_textbox(s, x + 0.1, y + 0.25, box_w - 0.2, 0.5,
                    name, size=12, bold=True, color=INK_1,
                    align=PP_ALIGN.CENTER)
        # arrow
        if i < 4:
            arrow = slide_arrow(s, x + box_w + 0.05, y + box_h / 2,
                                x + box_w + 0.35, y + box_h / 2)
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = INK_3
            arrow.line.fill.background()

    # Two-column detail
    col_w = (sw - 1.3) / 2
    top = 4.0
    col_h = 2.9
    add_rect(s, 0.55, top, col_w, col_h, BG_2)
    add_rect(s, 0.55, top, col_w, 0.06, ACCENT)
    add_textbox(s, 0.75, top + 0.18, col_w - 0.4, 0.4,
                "Features",
                size=14, bold=True, color=INK_1)
    add_bullets(s, 0.75, top + 0.6, col_w - 0.4, col_h - 0.7, [
        "Local Ollama integration (qwen3:4b, gpt-oss:120b-cloud, …)",
        "TF-IDF fallback when Ollama is offline",
        "Dataset-aware retrieval — facts walked from analysis result",
        "cosine similarity:  sim(q, cᵢ) = (q · cᵢ) / (‖q‖ ‖cᵢ‖)",
        "Sandboxed pandas via AST allow-list (safe_pandas.py)",
    ], size=12)

    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, top, col_w, col_h, BG_2)
    add_rect(s, rx, top, col_w, 0.06, ACCENT2)
    add_textbox(s, rx + 0.2, top + 0.18, col_w - 0.4, 0.4,
                "Why grounded?",
                size=14, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, top + 0.6, col_w - 0.4, col_h - 0.7, [
        "Every answer cites the analysis dictionary, not generated prose",
        "Planner → Executor → Critic → Writer (full path)",
        "Fast path — single writer call with brief + retrieved facts",
        "Result: AI answers are reliable, explainable, mathematically grounded.",
    ], size=12)
    return s


def slide_arrow(slide, x1, y1, x2, y2):
    """Draw a small horizontal arrow between two points (Inches)."""
    return slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x1), Inches(y1 - 0.12),
        Inches(x2 - x1), Inches(0.24),
    )


def slide_18_comparison(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Tool Comparison",
               subtitle="Where DataLens AI lands against the four families of competing tools.")

    headers = ["Feature", "DataLens AI", "Profiling Tools",
               "Cloud AutoML", "Conversational AI"]
    rows = [
        ("Offline Processing",     "✓",  "✓",     "—",     "—"),
        ("ML Leaderboard",         "✓",  "—",     "✓",     "Partial"),
        ("SHAP Explainability",    "✓",  "—",     "Limited", "—"),
        ("Drift Detection",        "✓",  "—",     "—",     "—"),
        ("Anomaly Ensemble",       "✓",  "—",     "—",     "—"),
        ("Time-Series Analytics",  "✓",  "—",     "Partial", "—"),
        ("AI Dataset Chat",        "✓",  "—",     "—",     "✓"),
        ("Privacy Friendly",       "✓",  "✓",     "—",     "—"),
    ]
    add_table(
        s,
        left=0.55, top=2.5, width=12.2, height=4.3,
        headers=headers, rows=rows,
        header_size=12, body_size=11,
        col_widths=[3.4, 1.8, 2.0, 2.0, 2.4],
        first_col_bold=True,
    )
    return s


def slide_19_limitations(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Limitations & Future Work",
               subtitle="What v3 doesn't solve yet — and where v4 should go.")
    sw = prs.slide_width / 914400
    col_w = (sw - 1.3) / 2

    add_rect(s, 0.55, 2.5, col_w, 4.4, BG_2)
    add_rect(s, 0.55, 2.5, col_w, 0.06, BAD)
    add_textbox(s, 0.75, 2.7, col_w - 0.4, 0.4,
                "Current limitations",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, col_w - 0.4, 3.5, [
        "Default ML hyperparameters — leaderboard is a baseline, not state-of-art",
        "Drift compares two halves of one upload (no persistent baseline)",
        "In-memory pipeline — degrades > 1 M rows",
        "Associational only — no causal inference",
        "English-centric NLP (NLTK stopwords, English seed lexicon)",
        "Snapshot-based — no real-time / streaming support",
    ], size=12)

    rx = 0.55 + col_w + 0.2
    add_rect(s, rx, 2.5, col_w, 4.4, BG_2)
    add_rect(s, rx, 2.5, col_w, 0.06, ACCENT)
    add_textbox(s, rx + 0.2, 2.7, col_w - 0.4, 0.4,
                "Future improvements",
                size=15, bold=True, color=INK_1)
    add_bullets(s, rx + 0.2, 3.15, col_w - 0.4, 3.5, [
        "Hyperparameter optimization with Optuna for the leaderboard",
        "Persistent drift baselines — store time-zero snapshot per dataset",
        "Streaming / sampled variants of deep stats and anomaly ensemble",
        "Causal inference (DoWhy / EconML) on declared treatment + outcome",
        "Multi-agent AI workflows (planner / executor / critic specialised)",
        "Multilingual NLP — language detection + per-language stopwords",
    ], size=12)
    return s


def slide_20_conclusion(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="Conclusion",
               subtitle="One unified offline platform — built end-to-end in one semester.")
    sw = prs.slide_width / 914400

    # Summary card
    add_rect(s, 0.55, 2.5, sw - 1.1, 2.0, BG_3)
    add_rect(s, 0.55, 2.5, sw - 1.1, 0.06, ACCENT)
    add_textbox(s, 0.75, 2.7, sw - 1.5, 0.4,
                "Summary",
                size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, 3.15, sw - 1.5, 1.3, [
        "DataLens AI integrates profiling, explainable ML, anomaly detection, time-series, drift, and conversational AI",
        "Every analytical claim traces back to a named algorithm and a concrete numeric result",
        "Offline-first by design  ·  privacy-friendly  ·  reproducible across runs (random_state=42)",
    ], size=12)

    # Final impact card
    add_rect(s, 0.55, 4.7, sw - 1.1, 2.2, BG_2)
    add_rect(s, 0.55, 4.7, sw - 1.1, 0.06, ACCENT2)
    add_textbox(s, 0.75, 4.9, sw - 1.5, 0.4,
                "Final impact",
                size=15, bold=True, color=INK_1)
    add_textbox(s, 0.75, 5.35, sw - 1.5, 1.4,
                "Transforms hours of manual data analysis into an automated, "
                "AI-driven workflow — without sending a single byte to the cloud.",
                size=15, color=INK_1)
    add_textbox(s, 0.75, 6.45, sw - 1.5, 0.35,
                "./run.sh  ·  17-tab dashboard at :5173  ·  console at :8501",
                size=11, color=INK_3, font="JetBrains Mono")
    return s


def slide_21_references(prs):
    s = add_blank_slide(prs)
    add_header(s, prs, title="References",
               subtitle="Algorithmic foundations and tooling cited throughout.")

    refs = [
        "Lundberg & Lee (2017). A Unified Approach to Interpreting Model Predictions. NeurIPS.",
        "Breiman (2001). Random Forests. Machine Learning, 45(1).",
        "Liu, Ting, Zhou (2008). Isolation Forest. ICDM.",
        "Breunig et al. (2000). LOF: Identifying Density-Based Local Outliers. SIGMOD.",
        "Cleveland et al. (1990). STL: A Seasonal-Trend Decomposition based on Loess.",
        "Hyndman & Athanasopoulos (2018). Forecasting: Principles and Practice (2nd ed.).",
        "Dickey & Fuller (1979). Distribution of the Estimators for AR Time Series with a Unit Root.",
        "Kwiatkowski, Phillips, Schmidt, Shin (1992). KPSS test for stationarity.",
        "Yurdakul (2018). Statistical Properties of Population Stability Index.",
        "Pedregosa et al. (2011). Scikit-learn: Machine Learning in Python. JMLR.",
        "Seabold & Perktold (2010). statsmodels: Econometric and Statistical Modeling with Python.",
        "Chen & Guestrin (2016). XGBoost: A Scalable Tree Boosting System. KDD.",
        "Ke et al. (2017). LightGBM: A Highly Efficient Gradient Boosting Decision Tree. NeurIPS.",
        "Han et al. (2024). PyOD 2: A Python Library for Outlier Detection. arXiv:2412.12154.",
        "Evidently AI Documentation  ·  PandasAI Documentation  ·  TableGPT / TableLLM papers.",
    ]
    sw = prs.slide_width / 914400
    add_bullets(s, 0.55, 2.5, sw - 1.1, 4.6, refs, size=12, color=INK_2)
    return s


def slide_22_thankyou(prs):
    s = add_blank_slide(prs)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400

    # Light hero background
    add_rect(s, 0, 0, sw, sh, BG_3)
    add_rect(s, 0, sh / 2 - 0.04, sw, 0.04, ACCENT)

    add_textbox(s, 0.6, sh / 2 - 1.8, sw - 1.2, 1.5,
                "Thank You.", size=72, bold=True, color=INK_1,
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.6, sh / 2 - 0.1, sw - 1.2, 0.6,
                "Questions?",
                size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.6, sh / 2 + 0.7, sw - 1.2, 0.4,
                "DATALENS AI  ·  GROUP 13",
                size=14, bold=True, color=INK_2,
                align=PP_ALIGN.CENTER, font="JetBrains Mono")
    add_textbox(s, 0.6, sh - 1.4, sw - 1.2, 0.4,
                "./run.sh  ·  17-tab dashboard at :5173  ·  console at :8501",
                size=11, color=INK_3,
                align=PP_ALIGN.CENTER, font="JetBrains Mono")
    return s


# ── Build driver ─────────────────────────────────────────────────────────────


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        slide_01_title,
        slide_02_contents,
        slide_03_introduction,
        slide_04_problem,
        slide_05_motivation,
        slide_06_literature,
        slide_07_gaps,
        slide_08_objectives,
        slide_09_architecture,
        slide_10_data_processing,
        slide_11_statistics,
        slide_12_anomaly,
        slide_13_ml_pipeline,
        slide_14_shap,
        slide_15_timeseries_drift,
        slide_16_reporting,
        slide_17_agentic,
        slide_18_comparison,
        slide_19_limitations,
        slide_20_conclusion,
        slide_21_references,
        slide_22_thankyou,
    ]

    for fn in builders:
        fn(prs)

    # Footers (skip cover and thank-you)
    slides = list(prs.slides)
    total = len(slides)
    for i, slide in enumerate(slides, start=1):
        if i in (1, total):
            continue
        add_footer(slide, prs, page_num=i, total=total)

    out_dir = ROOT / "reports"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "datalens_full_deck.pptx"
    prs.save(out_path)
    size_kb = out_path.stat().st_size / 1024
    print(f"  PPTX → {out_path}  ({total} slides, {size_kb:.0f} KB)")
    return out_path


if __name__ == "__main__":
    build()
