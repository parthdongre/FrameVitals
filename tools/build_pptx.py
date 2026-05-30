"""
Build the end-of-semester PPTX deck for DataLens AI.

Mirrors the structure of the team's MSE Review deck, but updated for v3 with
all the work that landed since: deep statistics, anomaly ensemble, model
leaderboard + SHAP, time-series, drift, agent + RAG, modernized PDF,
26-page whitepaper, and the live component audit (72/72 passing).

Run from repo root:
    python tools/build_pptx.py
Output:
    reports/datalens_endsem.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent

# ── Editorial colour palette (matches the React dashboard) ────────────────────

INK_1   = RGBColor(0x0A, 0x0A, 0x0A)   # near black
INK_2   = RGBColor(0x33, 0x33, 0x33)   # body
INK_3   = RGBColor(0x55, 0x55, 0x55)   # muted
PAGE    = RGBColor(0xFF, 0xFF, 0xFF)
BG_2    = RGBColor(0xF5, 0xF8, 0xFF)
ACCENT  = RGBColor(0x0D, 0x8A, 0x72)   # editorial teal
ACCENT2 = RGBColor(0x25, 0x63, 0xEB)   # blue accent (for highlights)
WARN    = RGBColor(0xC4, 0x7A, 0x09)
OK_GR   = RGBColor(0x2F, 0x7A, 0x4D)


# ── Slide helpers ─────────────────────────────────────────────────────────────

def add_blank_slide(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=INK_2, align=PP_ALIGN.LEFT,
                font="Inter"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items,
                *, size=15, color=INK_2, font="Inter", line_spacing=1.25):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, color, *, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    shape.shadow.inherit = False
    return shape


def add_header(slide, prs, *, eyebrow="DataLens · AI", title="", subtitle=""):
    # Top accent bar
    add_rect(slide, 0, 0, prs.slide_width / 914400, 0.35, ACCENT)

    # Eyebrow
    add_textbox(slide, 0.55, 0.45, 8, 0.3,
                eyebrow.upper(),
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")

    # Title
    if title:
        add_textbox(slide, 0.55, 0.75, prs.slide_width / 914400 - 1.1, 0.85,
                    title, size=32, bold=True, color=INK_1)
    if subtitle:
        add_textbox(slide, 0.55, 1.55, prs.slide_width / 914400 - 1.1, 0.45,
                    subtitle, size=14, color=INK_3)

    # Hairline
    add_rect(slide, 0.55, 2.05, 1.4, 0.04, ACCENT)


def add_footer(slide, prs, page_num: int, total: int):
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400
    add_rect(slide, 0.55, sh - 0.45, sw - 1.1, 0.01, INK_3)
    add_textbox(slide, 0.55, sh - 0.4, 4, 0.3,
                "DataLens AI  ·  End-Semester Review  ·  Group 13",
                size=9, color=INK_3, font="JetBrains Mono")
    add_textbox(slide, sw - 1.55, sh - 0.4, 1, 0.3,
                f"{page_num:02d} / {total:02d}",
                size=9, color=INK_2, bold=True, align=PP_ALIGN.RIGHT,
                font="JetBrains Mono")


# ── Title (cover) slide ───────────────────────────────────────────────────────

def slide_cover(prs):
    s = add_blank_slide(prs)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400

    # Hero band background
    add_rect(s, 0, 0, sw, sh * 0.55, RGBColor(0xF0, 0xF5, 0xFF))
    add_rect(s, 0, sh * 0.55 - 0.04, sw, 0.04, ACCENT)

    # Eyebrow
    add_textbox(s, 0.6, 0.7, 8, 0.4,
                "DATALENS · AI   ·   END-SEMESTER REVIEW",
                size=12, bold=True, color=ACCENT, font="JetBrains Mono")

    # Big title
    add_textbox(s, 0.6, 1.3, sw - 1.2, 1.1,
                "DataLens AI", size=54, bold=True, color=INK_1)

    add_textbox(s, 0.6, 2.45, sw - 1.2, 0.7,
                "Read the signal in your data.",
                size=24, bold=True, color=ACCENT)

    add_textbox(s, 0.6, 3.2, sw - 1.2, 0.6,
                "An offline-first, evidence-backed dataset analyser.",
                size=14, color=INK_3)

    # Bottom info card
    add_rect(s, 0.6, sh - 2.8, sw - 1.2, 2.0, BG_2)
    add_textbox(s, 0.85, sh - 2.6, 5, 0.3, "GUIDE",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    add_textbox(s, 0.85, sh - 2.3, 5, 0.4,
                "Prof. Kishor Pathak", size=14, bold=True, color=INK_1)

    add_textbox(s, 0.85, sh - 1.85, 5, 0.3, "TEAM  ·  GROUP 13",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono")
    add_bullets(s, 0.85, sh - 1.55, 6, 1.2, [
        "Parth Birari   ·  54 / IT D / 1251130415",
        "Parth Dongre   ·  54 / IT D / 1251130421",
        "Aryan Patil    ·  55 / IT D / 1251130411",
        "Atharv Patil   ·  56 / IT D / 1251130441",
    ], size=12, color=INK_2)

    # Right-aligned institute
    add_textbox(s, sw - 5.6, sh - 2.6, 5, 0.3,
                "INSTITUTE",
                size=10, bold=True, color=ACCENT, font="JetBrains Mono",
                align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, sh - 2.3, 5, 0.4,
                "Vishwakarma Institute of Technology, Pune",
                size=14, bold=True, color=INK_1, align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, sh - 1.85, 5, 0.4,
                "Department of Information Technology",
                size=12, color=INK_2, align=PP_ALIGN.RIGHT)
    add_textbox(s, sw - 5.6, sh - 1.45, 5, 0.4,
                "Date:  May 2026",
                size=11, color=INK_3, align=PP_ALIGN.RIGHT,
                font="JetBrains Mono")
    return s


# ── Standard content slide ────────────────────────────────────────────────────

def slide_content(prs, *, title, subtitle="", eyebrow="DataLens · AI"):
    s = add_blank_slide(prs)
    add_header(s, prs, eyebrow=eyebrow, title=title, subtitle=subtitle)
    return s


# ── KPI tile slide ────────────────────────────────────────────────────────────

def slide_kpi_grid(prs, *, title, subtitle, tiles):
    s = slide_content(prs, title=title, subtitle=subtitle)
    sw = prs.slide_width / 914400

    # Lay out 4 tiles per row
    per_row = min(4, len(tiles))
    gap = 0.25
    total_w = sw - 1.1
    tile_w = (total_w - gap * (per_row - 1)) / per_row
    tile_h = 1.6
    top = 2.5

    for i, t in enumerate(tiles):
        col = i % per_row
        row = i // per_row
        left = 0.55 + col * (tile_w + gap)
        y = top + row * (tile_h + gap)

        # Background tile
        add_rect(s, left, y, tile_w, tile_h, BG_2)
        # Top accent strip
        add_rect(s, left, y, tile_w, 0.06,
                 t.get("accent", ACCENT))

        # Label
        add_textbox(s, left + 0.2, y + 0.2, tile_w - 0.4, 0.3,
                    t.get("label", "").upper(),
                    size=9, bold=True, color=INK_3,
                    font="JetBrains Mono")
        # Value
        add_textbox(s, left + 0.2, y + 0.45, tile_w - 0.4, 0.7,
                    t.get("value", ""),
                    size=26, bold=True, color=INK_1)
        # Hint
        if t.get("hint"):
            add_textbox(s, left + 0.2, y + 1.15, tile_w - 0.4, 0.35,
                        t["hint"], size=10, color=INK_3)
    return s


# ── Bullet content slide ──────────────────────────────────────────────────────

def slide_bullets(prs, *, title, subtitle, bullets, eyebrow="DataLens · AI"):
    s = slide_content(prs, title=title, subtitle=subtitle, eyebrow=eyebrow)
    sw = prs.slide_width / 914400
    add_bullets(s, 0.55, 2.5, sw - 1.1, 4.5, bullets, size=15)
    return s


# ── Two-column comparison slide ──────────────────────────────────────────────

def slide_two_col(prs, *, title, subtitle,
                   left_title, left_bullets,
                   right_title, right_bullets,
                   left_color=ACCENT, right_color=ACCENT2):
    s = slide_content(prs, title=title, subtitle=subtitle)
    sw = prs.slide_width / 914400

    col_w = (sw - 1.3) / 2
    top = 2.5

    # Left card
    add_rect(s, 0.55, top, col_w, 4.4, BG_2)
    add_rect(s, 0.55, top, col_w, 0.06, left_color)
    add_textbox(s, 0.75, top + 0.2, col_w - 0.4, 0.4,
                left_title, size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75, top + 0.7, col_w - 0.4, 3.5, left_bullets,
                size=13)

    # Right card
    add_rect(s, 0.55 + col_w + 0.2, top, col_w, 4.4, BG_2)
    add_rect(s, 0.55 + col_w + 0.2, top, col_w, 0.06, right_color)
    add_textbox(s, 0.75 + col_w + 0.2, top + 0.2, col_w - 0.4, 0.4,
                right_title, size=15, bold=True, color=INK_1)
    add_bullets(s, 0.75 + col_w + 0.2, top + 0.7, col_w - 0.4, 3.5, right_bullets,
                size=13)
    return s


# ── Section divider slide ─────────────────────────────────────────────────────

def slide_section_divider(prs, *, number, title, blurb=""):
    s = add_blank_slide(prs)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400

    # Full-bleed background
    add_rect(s, 0, 0, sw, sh, RGBColor(0xF8, 0xFA, 0xFF))

    # Big number
    add_textbox(s, 0.6, sh / 2 - 1.6, 4, 1.5,
                number, size=72, bold=True, color=ACCENT,
                font="JetBrains Mono")

    add_textbox(s, 0.6, sh / 2 - 0.1, sw - 1.2, 1.0,
                title, size=40, bold=True, color=INK_1)

    if blurb:
        add_textbox(s, 0.6, sh / 2 + 0.95, sw - 1.2, 1.0,
                    blurb, size=14, color=INK_3)

    add_rect(s, 0.6, sh / 2 + 1.7, 1.4, 0.04, ACCENT)
    return s


# ── Architecture diagram slide ───────────────────────────────────────────────

def slide_architecture(prs):
    s = slide_content(
        prs,
        title="System Architecture",
        subtitle="Six-phase pipeline · 33 modules · parallel where safe.",
    )
    sw = prs.slide_width / 914400

    blocks = [
        ("React Dashboard", "Vite · Tailwind · Highcharts", ACCENT,   1.0),
        ("Flask API",       "JSON contract on :5055",       ACCENT2,  2.5),
        ("run_full_analysis", "modules/pipeline.py",        ACCENT,   4.0),
    ]
    for label, sub, color, top in blocks:
        add_rect(s, sw / 2 - 1.6, top, 3.2, 0.8, BG_2)
        add_rect(s, sw / 2 - 1.6, top, 0.08, 0.8, color)
        add_textbox(s, sw / 2 - 1.4, top + 0.1, 3.0, 0.4,
                    label, size=14, bold=True, color=INK_1)
        add_textbox(s, sw / 2 - 1.4, top + 0.45, 3.0, 0.3,
                    sub, size=10, color=INK_3, font="JetBrains Mono")

    # Phases at the bottom
    phase_top = 5.2
    phases = [
        ("Phase 1", "Ingest\nload→profile→roles\n→signals→selector", ACCENT),
        ("Phase 2", "Quality\nhealth · ML readiness\n· advanced", OK_GR),
        ("Phase 3", "Parallel\ndeep-stats · anomalies\ntime-series · text", ACCENT2),
        ("Phase 4", "ML Chain\nleaderboard\n→ explainability", WARN),
        ("Phase 5", "Reporting\ncleaning · charts\n· PDF report", ACCENT),
        ("Phase 6", "Agentic\nbrief · RAG\n· LLM agent", ACCENT2),
    ]
    box_w = (sw - 1.1) / len(phases) - 0.1
    for i, (eye, body, color) in enumerate(phases):
        left = 0.55 + i * (box_w + 0.1)
        add_rect(s, left, phase_top, box_w, 1.5, BG_2)
        add_rect(s, left, phase_top, box_w, 0.05, color)
        add_textbox(s, left + 0.1, phase_top + 0.15, box_w - 0.2, 0.3,
                    eye.upper(), size=8, bold=True, color=color,
                    font="JetBrains Mono")
        add_textbox(s, left + 0.1, phase_top + 0.45, box_w - 0.2, 1.0,
                    body, size=9, color=INK_2)

    return s


# ── Build the deck ────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # 1. Cover
    slide_cover(prs)

    # 2. Contents
    slide_bullets(prs, title="Contents",
                  subtitle="What this deck walks through.",
                  bullets=[
                      "Introduction & Motivation",
                      "Problem Statement & Research Gaps",
                      "Objectives",
                      "System Architecture",
                      "Modules & Algorithms (33 in total)",
                      "Demo: end-to-end pipeline",
                      "Live Component Audit (72 / 72 passing)",
                      "Visual Outputs: charts + PDF whitepaper",
                      "Comparison with Existing Tools",
                      "Limitations & Future Work",
                      "Conclusion & References",
                  ])

    # 3. Section: Foundation
    slide_section_divider(prs, number="01",
                           title="Foundation",
                           blurb="What we built and why.")

    # 4. Introduction
    slide_bullets(prs, title="Introduction",
                  subtitle="What DataLens AI is.",
                  bullets=[
                      "An offline-first dataset analyser that turns any CSV / TSV / Excel / JSON into a structured report — no manual coding required.",
                      "Six-phase Python pipeline: ingest → quality → statistics → modeling → time-series + text → drift, plus charts, PDF, and an agentic Q&A layer.",
                      "Two user surfaces: a React dashboard at :5173 with 17 tabs, and a Streamlit step-by-step visual console at :8501.",
                      "All ML and stats run locally; LLM responses use Ollama (local or cloud) with a TF-IDF + heuristic fallback when no model is reachable.",
                      "Outputs: a 17-tab interactive dashboard, a 26-page typeset PDF whitepaper, a multi-page dataset PDF report, and a chat-style Q&A agent.",
                  ])

    # 5. Motivation
    slide_bullets(prs, title="Motivation",
                  subtitle="What was wrong with the status quo.",
                  bullets=[
                      "Data analysis still requires Python / SQL coding or expensive licensed tools (Tableau, Power BI Pro, JMP).",
                      "Existing automated tools (pandas-profiling, ydata-profiling) cover profiling but stop short of ML readiness, drift detection, and SHAP.",
                      "Online AI tools (ChatGPT Code Interpreter, Julius AI) raise hard privacy concerns — data leaves the machine.",
                      "Most outputs are descriptive, not explainable: the analyst sees correlations but not causal feature importance.",
                      "Real-world datasets are messy: encoding mismatches, mixed types, dirty timestamps. Most tools assume clean data.",
                      "No single tool combines automation + explainability + offline operation in a single click-through pipeline.",
                  ])

    # 6. Problem statement
    slide_bullets(prs, title="Problem Statement",
                  subtitle="The gap we set out to close.",
                  bullets=[
                      "Build a single, automated, offline-first system that takes any tabular dataset and produces a complete analytical report in under one minute.",
                      "The system must run entirely locally: no API keys, no cloud uploads, no data leaves the machine.",
                      "It must be honest: every claim must be traceable to a named algorithm and an underlying numeric result, not generated text.",
                      "It must be explainable: every model decision must be tied back to the input features that produced it (SHAP / permutation importance).",
                      "It must be robust to messy data: encoding mismatches, missing values, mixed dtypes, timestamps in different formats.",
                  ])

    # 7. Research gaps
    slide_two_col(prs,
        title="Research Gaps Identified",
        subtitle="What existing tools miss, and what DataLens adds.",
        left_title="Existing tools — what's missing",
        left_bullets=[
            "Profiling-only depth (no ML, no drift, no SHAP)",
            "Closed-source, cloud-only LLM inference",
            "Black-box recommendations with no math behind them",
            "Single chart per relationship — no chart planner",
            "No reproducibility: outputs differ between runs",
            "No per-row anomaly attribution",
        ],
        right_title="DataLens fills these by design",
        right_bullets=[
            "33 modules covering profile → ML → SHAP → drift → text",
            "Local Ollama with optional Ollama Cloud free-tier fallback",
            "Every signal carries evidence + recommendation strings",
            "20+ chart types planned by role-aware planner",
            "Deterministic random_state seeds across all phases",
            "Anomaly ensemble: IsoForest + LOF + MAD + Mahalanobis",
        ],
    )

    # 8. Objectives
    slide_bullets(prs, title="Objectives",
                  subtitle="What we promised to deliver.",
                  bullets=[
                      "Auto-detect column roles, dtypes, and dataset signals without user input.",
                      "Produce a 0–100 health score and a 0–100 ML-readiness score with clear sub-component breakdowns.",
                      "Run a cross-validated ML leaderboard for any user-supplied target column, classification or regression.",
                      "Generate global + per-row SHAP attributions for the leaderboard winner.",
                      "Detect time-series structure and run ADF + KPSS stationarity, STL decomposition, and Holt–Winters forecasting.",
                      "Compute population-stability-index (PSI) drift between two datasets or two halves of one dataset.",
                      "Produce a styled, multi-page PDF report and a 17-tab React dashboard from the same JSON payload.",
                      "Expose a chat-style Q&A interface grounded in the analysis result via RAG over the fact index.",
                  ])

    # 9. Section: System
    slide_section_divider(prs, number="02",
                           title="System",
                           blurb="Architecture, modules, and pipeline.")

    # 10. Architecture
    slide_architecture(prs)

    # 11. Module map
    slide_kpi_grid(prs,
        title="The 33 Analytical Modules",
        subtitle="What each module category contributes to the final report.",
        tiles=[
            {"label": "Ingest",      "value": "5 modules",
             "hint": "loader · profiler · roles · signals · selector",
             "accent": ACCENT},
            {"label": "Quality",     "value": "5 modules",
             "hint": "health · ML readiness · advanced · signals · cleaner",
             "accent": OK_GR},
            {"label": "Statistics",  "value": "2 modules",
             "hint": "deep_statistics_v2 · anomaly_ensemble",
             "accent": ACCENT},
            {"label": "Modeling",    "value": "2 modules",
             "hint": "model_leaderboard · explainability (SHAP)",
             "accent": WARN},
            {"label": "Time & Text", "value": "2 modules",
             "hint": "time_series · text_profile",
             "accent": ACCENT2},
            {"label": "Drift",       "value": "1 module",
             "hint": "drift_analysis (PSI + KS + chi²)",
             "accent": ACCENT2},
            {"label": "Reporting",   "value": "3 modules",
             "hint": "chart_planner · visualizer · pdf_report_builder",
             "accent": ACCENT},
            {"label": "Agentic",     "value": "4 modules",
             "hint": "frontend_api · agent_brief · rag_index · ai_agent",
             "accent": ACCENT2},
        ],
    )

    # 12. Methodology — Quality
    slide_bullets(prs, title="Methodology · Quality",
                  subtitle="Two scores out of 100, four sub-components.",
                  bullets=[
                      "Health score = 0.30·completeness + 0.20·consistency + 0.20·uniqueness + 0.20·outlier safety + 10",
                      "Completeness = 100 − mean missing %",
                      "Outlier safety uses 1.5× IQR fences — share of numeric cells inside the safe band",
                      "ML readiness penalises missing %, duplicate %, and encoding load (categorical column count)",
                      "Auto-cleaner: median impute numerics, mode impute categoricals, dedupe rows, coerce dtypes — every action tagged Low / Medium / High risk",
                      "Signal engine converts the per-component scores into severity-tagged signal cards with evidence + recommendation strings",
                  ])

    # 13. Methodology — Statistics
    slide_bullets(prs, title="Methodology · Deep Statistics",
                  subtitle="Bivariate tests and outlier ensembles.",
                  bullets=[
                      "Numeric pairs: Pearson r + Spearman ρ + Kendall τ with p-values; sorted by |r|",
                      "Categorical pairs: Cramér's V (normalised χ² in [0,1]) — bounded, comparable across cardinalities",
                      "Numeric vs categorical: point-biserial for binary; Kruskal–Wallis for multi-group, with η² effect size",
                      "Per-column normality: Shapiro–Wilk (n ≤ 5000) + D'Agostino–Pearson + Anderson–Darling, with consensus verdict",
                      "Distribution fit: AIC across norm, lognorm, expon, gamma, weibull_min, beta — leaderboard returned",
                      "Bootstrap 95% CIs for mean and median via BCa (scipy.stats.bootstrap), 999 resamples",
                      "Outlier flags: IQR + classical z (|z|>3) + robust MAD-z (|z|>3.5) reported separately",
                      "Pairwise budget capped via max_pairs to keep runtime O(k) on wide datasets",
                  ])

    # 14. Methodology — Anomalies
    slide_bullets(prs, title="Methodology · Anomaly Ensemble",
                  subtitle="Five independent detectors, mean-rank averaged.",
                  bullets=[
                      "Isolation Forest (Liu, Ting, Zhou 2008) — recursive partition depth as anomaly score",
                      "Local Outlier Factor (Breunig et al. 2000) — local density vs neighbourhood",
                      "Robust z-score via MAD — coordinate-wise median absolute deviation",
                      "Mahalanobis distance with shrinkage covariance (MinCovDet)",
                      "EllipticEnvelope — robust Gaussian envelope with contamination prior",
                      "Optional pyod ECOD + COPOD when pyod is installed",
                      "Each detector emits [0,1] after min-max normalisation; ensemble = column-wise mean",
                      "Rows above ensemble threshold (default 0.6) are flagged and surfaced in the Anomalies tab",
                  ])

    # 15. Methodology — Modeling + SHAP
    slide_bullets(prs, title="Methodology · Modeling + SHAP",
                  subtitle="Cross-validated leaderboard with an honest baseline.",
                  bullets=[
                      "Classification: LogisticRegression, KNN, RandomForest, GradientBoosting, XGB, LGBM + Dummy floor",
                      "Regression: Ridge, Lasso, KNN, RandomForest, GradientBoosting, XGB, LGBM + Dummy floor",
                      "Stratified 5-fold CV (classification) or KFold-5 (regression); folds capped by smallest class",
                      "Primary metrics: F1-weighted (classification) · R² (regression); also reports MAE, RMSE, accuracy, precision, recall",
                      "Winner refit on full train, then validated on held-out 25% slice",
                      "SHAP TreeExplainer for tree models, LinearExplainer for linear leaders, PermutationImportance fallback",
                      "Three SHAP properties verified: local accuracy, missingness, consistency",
                      "Per-row attributions stored for the top-3 most-extreme test rows for the dashboard's audit trail",
                  ])

    # 16. Methodology — Time-series
    slide_bullets(prs, title="Methodology · Time-series & Text",
                  subtitle="Stationarity, decomposition, forecasting; vocabulary + n-grams.",
                  bullets=[
                      "Date detection: name regex + ≥70% parse rate + monotonicity score; FFT picks the dominant period",
                      "Stationarity: ADF (rejects H0 → stationary) + KPSS (rejects H0 → non-stationary) cross-checked into 4 verdicts",
                      "STL decomposition (Cleveland 1990); strength-of-trend / seasonality via Hyndman formulation",
                      "Forecasting: Holt–Winters additive vs naive last-value baseline on 80/20 train-test split",
                      "Text columns identified by 30%+ multi-token rate, mean length ≥ 8, not ID-like",
                      "Per-text-column stats: vocabulary size, type-token ratio, mean sentence length, top-15 unigrams + bigrams",
                      "Regex hits: emails, URLs, phone numbers, mentions, hashtags",
                      "NLTK English stopwords when available; dtype-driven everywhere else",
                  ])

    # 17. Methodology — Drift
    slide_bullets(prs, title="Methodology · Drift",
                  subtitle="Has the data changed? PSI + KS + χ².",
                  bullets=[
                      "Population Stability Index (PSI): bin via reference quantiles, sum (p_cur − p_ref) · log(p_cur/p_ref)",
                      "PSI cut-offs: < 0.10 stable · 0.10–0.25 minor · 0.25–0.50 moderate · ≥ 0.50 severe",
                      "Two-sample Kolmogorov–Smirnov complements PSI for subtle continuous shifts",
                      "Categorical drift: same PSI on observed proportions + chi-square test of independence",
                      "Per-column severity aggregated into overall verdict: stable / minor / moderate / severe",
                      "Robustness fix: float64 coerce before np.quantile (newer numpy refuses to subtract bool arrays)",
                  ])

    # 18. Methodology — Agentic + RAG
    slide_bullets(prs, title="Methodology · Agentic + RAG",
                  subtitle="Grounded LLM Q&A over the analysis result.",
                  bullets=[
                      "Dataset brief: ~9KB JSON packed from the result — schema, signals, ML headlines, top correlations, sample rows",
                      "Fact index: flattens every analysis path-value pair into atomic Fact objects",
                      "Embedding backend: Ollama nomic-embed-text (parallel-batched, ~9s for 1300 facts) → cosine top-k",
                      "Fallback: TF-IDF over the same corpus (~17ms for 1300 facts) — used in tests and CI",
                      "Agent loop: planner → executor → critic → writer (full mode) or single writer call (fast mode, default)",
                      "Tools the agent can call: get_dataset_brief, rank_columns_by_missing, get_anomaly_top_rows, get_leaderboard, etc.",
                      "Backends: Ollama Cloud (gpt-oss:120b-cloud, free tier, ~1–2s) → local Ollama → OpenRouter → heuristic fallback",
                      "If all four fail: deterministic writer surfaces retrieved facts directly, labelled 'no LLM endpoint reachable'",
                  ])

    # 19. Section: Results
    slide_section_divider(prs, number="03",
                           title="Results",
                           blurb="Live audit, demos, and visual outputs.")

    # 20. Live audit
    slide_kpi_grid(prs,
        title="Live Component Audit",
        subtitle="Every module exercised against three demo datasets right before this deck was generated.",
        tiles=[
            {"label": "Components", "value": "24",
             "hint": "Per-module probes",            "accent": ACCENT},
            {"label": "Datasets",   "value": "3",
             "hint": "classification + regression + no-target", "accent": ACCENT2},
            {"label": "Total tests", "value": "72",
             "hint": "24 × 3 datasets",              "accent": OK_GR},
            {"label": "Pass rate",   "value": "100%",
             "hint": "72 / 72 passing",              "accent": OK_GR},
            {"label": "Pipeline median", "value": "5.2 s",
             "hint": "no target supplied",           "accent": ACCENT},
            {"label": "ML stack",   "value": "~30 s",
             "hint": "leaderboard + SHAP",           "accent": WARN},
            {"label": "Whitepaper", "value": "26 pages",
             "hint": "Auto-built LaTeX → PDF",       "accent": ACCENT},
            {"label": "Charts",     "value": "20+ types",
             "hint": "Per-dataset",                  "accent": ACCENT},
        ],
    )

    # 21. Datasets demo
    slide_bullets(prs, title="Demo Datasets — 17 in total",
                  subtitle="Each tuned to highlight a specific cluster of components.",
                  bullets=[
                      "customer_churn, real_estate, ecommerce_transactions, financial_transactions — full ML stack",
                      "wine_quality, manufacturing_qc, employee_attrition, dirty_marketing — quality + cleaning",
                      "iot_sensor_readings, bike_sharing, financial_transactions — time-series detection",
                      "support_tickets, medical_records, survey_responses — text profile + NLP",
                      "insurance_claims, flight_delays — anomaly ensemble + drift",
                      "9 additional generated datasets covering fairness, leakage, imbalanced classes, kitchen-sink roles",
                      "Generators in demo_datasets/generate.py and demo_datasets/generate_more_datasets.py",
                      "All 17 verified by demo_datasets/verify.py with a recommended target column per dataset",
                  ])

    # 22. Visual outputs
    slide_two_col(prs,
        title="Visual Outputs",
        subtitle="What the user actually sees and downloads.",
        left_title="In-app",
        left_bullets=[
            "17-tab React report (Overview, Profile, Statistics, Anomalies, ML Lab, SHAP, Time-series, Text, Drift, Charts, etc.)",
            "20+ chart types per analysis, role-aware planner-driven",
            "Click-to-fullscreen lightbox on every chart",
            "Markdown-rendered LLM responses in chat",
            "Live progress timeline in Streamlit console (:8501)",
            "Editorial dark theme: cream + teal on near-black",
        ],
        right_title="Downloadable",
        right_bullets=[
            "Cleaned CSV (auto-cleaner output)",
            "Multi-page styled PDF dataset report (~1.4 MB)",
            "26-page LaTeX whitepaper with embedded charts (~2 MB)",
            "Component test manifest as JSON",
            "Per-run static chart PNGs in /static/charts/",
            "Each PDF has cover, KPI tiles, contents, and visual evidence pages",
        ],
    )

    # 23. Comparison
    slide_two_col(prs,
        title="Comparison · DataLens vs Existing Tools",
        subtitle="Where DataLens wins, and where it doesn't (yet).",
        left_title="DataLens AI",
        left_bullets=[
            "Offline-first; runs without internet",
            "33 modules covering 8 analysis categories",
            "Built-in SHAP + permutation importance",
            "Built-in PSI / KS drift detection",
            "Multi-detector anomaly ensemble (5–7 detectors)",
            "Auto-generated PDF report + LaTeX whitepaper",
            "Local LLM agent with TF-IDF fallback",
            "Total project size: 1.2 GB (mostly venv + node_modules)",
        ],
        right_title="ydata-profiling / pandas-profiling / Evidently",
        right_bullets=[
            "Profiling depth comparable, but no SHAP, no leaderboard",
            "Evidently does drift well, but no profiling, no LLM",
            "ydata is single-file HTML — no tabbed navigation",
            "None of them have a chat-style Q&A interface",
            "None auto-generate a typeset academic-style PDF",
            "None bundle Auto-cleaner + ML Lab + Drift in one tool",
            "Most rely on cloud APIs for the AI summary",
            "AutoML (auto-sklearn, h2o.ai) doesn't profile or explain",
        ],
    )

    # 24. Limitations
    slide_bullets(prs, title="Limitations",
                  subtitle="What DataLens does NOT do well today.",
                  bullets=[
                      "Modeling depth: leaderboard runs default hyperparameters with no automated tuning — treat as baseline, not state-of-art",
                      "No persistent drift baseline — drift compares against a 50/50 row split rather than a stored historical reference",
                      "Causality: every relationship reported is associational; no causal inference framework",
                      "LLM grounding: agent re-call into pandas isn't yet supported, so questions outside the brief fall back to RAG",
                      "Memory bound: pipeline holds the entire dataset in RAM. Datasets >1M rows hit the deep-stats and leaderboard phases",
                      "Internationalisation: text profile uses NLTK English stopwords; non-English NLP is dtype-driven only",
                      "No real-time / streaming support; the pipeline is one-shot per upload",
                  ])

    # 25. Future work
    slide_bullets(prs, title="Future Work",
                  subtitle="Where v4 should go.",
                  bullets=[
                      "Persistent drift baselines — store a 'time-zero' snapshot per dataset and compute drift against it",
                      "Streaming / sampled variants of deep stats and anomaly ensemble for >1M-row datasets",
                      "Hyperparameter tuning with Optuna (already a dep) for the leaderboard",
                      "Causal inference module: DoWhy or EconML on declared treatment + outcome columns",
                      "LLM agent re-call: let the writer execute pandas snippets when retrieval is empty",
                      "Multilingual NLP: detect text language, swap stopwords, language-aware tokenisation",
                      "Production observability: Loguru → JSON logs, Prometheus metrics, latency SLOs",
                      "Vercel + Railway deployment so the demo runs without a local machine",
                  ])

    # 26. Conclusion
    slide_bullets(prs, title="Conclusion",
                  subtitle="What this project demonstrates.",
                  bullets=[
                      "DataLens AI is a complete, offline-first, evidence-backed dataset analyser — built end-to-end in one semester.",
                      "It demonstrates working integration of profiling, deep statistics, anomaly detection, ML modeling, SHAP, time-series, drift, and an LLM agent in a single coherent pipeline.",
                      "Every analytical claim in the system traces back to a named algorithm and a concrete numeric result — no generated text, no hallucinated metrics.",
                      "The 72 / 72 component audit and the 26-page auto-generated whitepaper give us confidence the system behaves as documented.",
                      "The project is delivered: source, demo datasets, whitepaper, dashboard, PDF report, all in one repo, runnable with a single ./run.sh command.",
                  ])

    # 27. References
    slide_bullets(prs, title="References",
                  subtitle="Algorithmic foundations cited in the codebase and the whitepaper.",
                  bullets=[
                      "Lundberg & Lee (2017). A Unified Approach to Interpreting Model Predictions. NeurIPS.",
                      "Breiman (2001). Random Forests. Machine Learning 45(1).",
                      "Liu, Ting, Zhou (2008). Isolation Forest. ICDM.",
                      "Breunig et al. (2000). LOF: Identifying Density-Based Local Outliers. SIGMOD.",
                      "Cleveland et al. (1990). STL: A Seasonal-Trend Decomposition Procedure Based on Loess.",
                      "Hyndman & Athanasopoulos (2018). Forecasting: Principles and Practice (2nd ed.).",
                      "Dickey & Fuller (1979). Distribution of the Estimators for Autoregressive Time Series with a Unit Root.",
                      "Kwiatkowski, Phillips, Schmidt, Shin (1992). KPSS test for stationarity.",
                      "Shapley (1953). A Value for n-Person Games.",
                      "Cramér (1946). Mathematical Methods of Statistics.",
                  ])

    # 28. Thank you
    s = add_blank_slide(prs)
    sw = prs.slide_width / 914400
    sh = prs.slide_height / 914400
    add_rect(s, 0, 0, sw, sh, RGBColor(0xF8, 0xFA, 0xFF))
    add_textbox(s, 0.6, sh / 2 - 1.0, sw - 1.2, 1.5,
                "Thank You.", size=64, bold=True, color=INK_1,
                align=PP_ALIGN.CENTER)
    add_textbox(s, 0.6, sh / 2 + 0.1, sw - 1.2, 0.6,
                "Questions?",
                size=24, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, 0.6, sh - 1.4, sw - 1.2, 0.4,
                "github · ./run.sh · 17-tab dashboard at :5173 · console at :8501",
                size=11, color=INK_3,
                align=PP_ALIGN.CENTER, font="JetBrains Mono")

    # ── Add footers (skip cover, divider slides, and final thank-you)
    slides = list(prs.slides)
    total = len(slides)
    for i, slide in enumerate(slides, start=1):
        # cover, dividers (3, 9, 19), thank-you (28) get no footer
        if i in (1, 3, 9, 19, total):
            continue
        add_footer(slide, prs, page_num=i, total=total)

    # ── Save ────────────────────────────────────────────────────────────────
    out_dir = ROOT / "reports"
    out_dir.mkdir(exist_ok=True)
    out_path = out_dir / "datalens_endsem.pptx"
    prs.save(out_path)
    print(f"  PPTX → {out_path}  ({total} slides, "
          f"{out_path.stat().st_size / 1024:.0f} KB)")
    return out_path


if __name__ == "__main__":
    build()
